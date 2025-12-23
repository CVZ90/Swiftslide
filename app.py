import streamlit as st
from groq import Groq
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import requests
import urllib.parse

# --- إعدادات الواجهة الاحترافية ---
st.set_page_config(page_title="SwiftSlide AI | VIP Edition", page_icon="💎")

st.markdown("""
    <style>
    .stApp { background-color: #0e1117; }
    .title-text { text-align: center; color: #00d4ff; font-size: 45px; font-weight: bold; text-shadow: 2px 2px #000; }
    .unlock-box { background: linear-gradient(145deg, #1e293b, #0f172a); padding: 30px; border-radius: 20px; border: 1px solid #3b82f6; text-align: center; box-shadow: 0 10px 15px rgba(0,0,0,0.5); }
    .whatsapp-btn { background-color: #25D366; color: white !important; padding: 12px 25px; border-radius: 30px; text-decoration: none; font-weight: bold; display: inline-block; margin-top: 15px; }
    </style>
""", unsafe_allow_html=True)

# --- دالة صناعة البوربوينت الاحترافي ---
def create_luxury_pptx(text_content, topic_name):
    prs = Presentation()
    
    # 1. سلايد العنوان الرئيسي
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(10, 25, 47) # كحلي داكن جداً
    
    title_box = slide.shapes.add_textbox(Inches(0), Inches(3), Inches(10), Inches(2))
    p = title_box.text_frame.add_paragraph()
    p.text = str(topic_name).upper()
    p.font.bold = True
    p.font.size = Pt(54)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 2. السلايدات الداخلية بتنسيق متناوب
    sections = str(text_content).split("Slide")
    for i, section in enumerate(sections):
        if len(section.strip()) > 15:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(15, 15, 15) # أسود فخم
            
            lines = section.strip().split("\n")
            title_text = str(lines[0]).replace(":", "").replace("1","").replace("2","").replace("3","").replace("4","").replace("5","").strip()
            
            # تحديد المواقع (تناوب بين اليمين واليسار)
            is_even = i % 2 == 0
            img_x = Inches(5.5) if is_even else Inches(0.5)
            text_x = Inches(0.5) if is_even else Inches(4.5)

            # إضافة العنوان مع خط زرق تحتيه
            t_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
            tp = t_box.text_frame.add_paragraph()
            tp.text = title_text
            tp.font.size = Pt(32)
            tp.font.bold = True
            tp.font.color.rgb = RGBColor(0, 212, 255)

            # إضافة صورة ذكية بناءً على عنوان السلايد
            try:
                search_query = title_text.split()[-1] # البحث بآخر كلمة في العنوان لضمان التنوع
                img_url = f"https://images.unsplash.com/photo-1550751827-4bd374c3f58b?q=80&w=400" # رابط احتياطي فخم
                # يمكنك استخدام: f"https://source.unsplash.com/400x300/?{search_query},technology"
                response = requests.get(img_url, timeout=5)
                slide.shapes.add_picture(io.BytesIO(response.content), img_x, Inches(1.5), width=Inches(4))
            except: pass

            # إضافة النقاط بتنسيق مريح
            b_box = slide.shapes.add_textbox(text_x, Inches(1.5), Inches(5), Inches(5))
            tf = b_box.text_frame
            tf.word_wrap = True
            for line in lines[1:]:
                if line.strip():
                    lp = tf.add_paragraph()
                    lp.text = "✦ " + str(line).strip("-*• ")
                    lp.font.size = Pt(18)
                    lp.font.color.rgb = RGBColor(220, 220, 220)
                    lp.space_after = Pt(12)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()

# --- المنطق الأساسي للتطبيق ---
st.markdown('<h1 class="title-text">SwiftSlide AI 🚀</h1>', unsafe_allow_html=True)
st.divider()

# سحب المفتاح من Secrets
api_key = st.secrets.get("GROQ_API_KEY")
if not api_key:
    st.error("Missing GROQ_API_KEY! Please add it to Streamlit Secrets.")
    st.stop()

client = Groq(api_key=api_key)

topic = st.text_input("Enter your topic for a VIP Presentation:", placeholder="e.g. The Future of AI in 2025")

if st.button("Generate Professional Presentation ✨"):
    if topic:
        with st.spinner("💎 AI is designing your premium slides..."):
            try:
                response = client.chat.completions.create(
                    model="llama-3.3-70b-versatile",
                    messages=[
                        {"role": "system", "content": "You are a professional presentation designer. Create 5 slides. For each slide, start with 'Slide X: Title' then 4 bullet points. English only."},
                        {"role": "user", "content": f"Topic: {topic}"}
                    ]
                )
                content = response.choices[0].message.content
                st.session_state['pptx_file'] = create_luxury_pptx(content, topic)
                st.session_state['topic'] = topic
                st.session_state['is_ready'] = True
                st.success("✅ Your presentation is ready to be unlocked!")
            except Exception as e:
                st.error(f"Error: {e}")

# --- منطقة الدفع والقفل ---
if st.session_state.get('is_ready'):
    st.markdown('<div class="unlock-box">', unsafe_allow_html=True)
    st.markdown(f"### 🔒 Unlock '{st.session_state['topic']}'")
    st.write("To download the High-Quality PPTX file, please send **$4** via **Whish Money** to:")
    st.code("81950506", language="text")
    
    code_input = st.text_input("Enter Activation Code:", type="password", key="active_key")
    
    if code_input == "SWIFT2025":
        st.balloons()
        st.download_button(
            label="📥 Download Luxury PPTX File",
            data=st.session_state['pptx_file'],
            file_name=f"{st.session_state['topic']}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )
    elif code_input != "":
        st.error("Invalid Code. Please contact support.")

    # زر الواتساب
    msg = f"Hello! I paid $4 for the presentation: {st.session_state['topic']}. Please send the code."
    wa_link = f"https://wa.me/96181950506?text={urllib.parse.quote(msg)}"
    st.markdown(f'<a href="{wa_link}" class="whatsapp-btn">I have paid - Send Proof</a>', unsafe_allow_html=True)
    st.markdown('</div>', unsafe_allow_html=True)

st.markdown("<br><p style='text-align: center; color: #4b5563;'>SwiftSlide AI Lebanon © 2025</p>", unsafe_allow_html=True)
